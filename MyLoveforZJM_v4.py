###
# 这是一个帮助当钉妹实现听写PPT自动化制作的小工具
# 要求
# 1. 需要把EXCEl文件准备好，这个陈乃平已经准备好了
# 2. 需要把对陈乃平爱准备好
# ###

import pandas as pd
import pptx as pptx
import re
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.dml.color import RGBColor

def parse_recite(strings):
    """
    解析词汇字符串，去掉数字和音标，并将英文单词和中文解释分别存放在两个数组中。

    :param strings: 包含词汇的字符串列表
    :return: 两个列表，分别包含英文单词和中文解释
    """
    # 初始化两个数组
    english_words = []
    chinese_translations = []

    # 定义正则表达式模式
    pattern_remove_number = re.compile(r'^\d+\.\s*')
    pattern_split_pronunciation = re.compile(r'\s*\[.*?\]\s*')
    # pattern_split_word_type = re.compile(r', (n\.|a\.|v\.|ad\.|conj\.|prep\.|pron\.|num\.|int\.|art\.|aux\.|modal\.|interj\.)')

    for line in strings:
        # 去掉数字和点号
        if not isinstance(line, str):
            continue  # 跳过非字符串类型的数据
        line = pattern_remove_number.sub('', line)
        
        # 去掉音标
        parts = pattern_split_pronunciation.split(line,maxsplit=1)
        
        # 提取英文单词和中文解释
        # parts = pattern_split_word_type.split(line, maxsplit=1)
        if len(parts) == 2:
            english_word = parts[0].strip()
            chinese_translation = parts[1].strip()
            chinese_translations.append(f"{chinese_translation}")
            english_words.append(english_word)
        else:
            raise ValueError(f"无法正确分割字符串: {line}")

    return english_words, chinese_translations

def parse_phrase(strings):
    """
    解析词汇字符串，去掉数字和点号，并将英文单词和中文解释分别存放在两个数组中。

    :param strings: 包含词汇的字符串列表
    :return: 两个列表，分别包含英文单词和中文解释
    """
    # 初始化两个数组
    english_words = []
    chinese_translations = []

    # 定义正则表达式模式
    pattern_remove_number = re.compile(r'^\d+\.\s*')
    pattern_split_english_chinese = re.compile(r'\s+(?=[\u4e00-\u9fff\u2026\uff08])')

    for line in strings:
        # 去掉数字和点号
        if not isinstance(line, str):
            continue  # 跳过非字符串类型的数据

        line = pattern_remove_number.sub('', line)
        
        # 按照英文和中文的分界线分割字符串
        parts = pattern_split_english_chinese.split(line, maxsplit=1)
        
        if len(parts) == 2:
            english_word = parts[0].strip()
            chinese_translation = parts[1].strip()
            english_words.append(english_word)
            chinese_translations.append(chinese_translation)
        else:
            raise ValueError(f"无法正确分割字符串: {line}")

    return english_words, chinese_translations



def sparse_transformation(strings):
# 输入字符串
    results_E = []
    results_C = []
    pattern_remove_number = re.compile(r'^\d+\.\s*')
    for string in strings:
        if not isinstance(string, str):
            continue  # 跳过非字符串类型的数据

        string = pattern_remove_number.sub('', string)
        parts = re.split(r'→', string)

        english_parts = [part.split(' ')[0] for part in parts]
        chinese_parts = [part[len(english_parts[i]):].strip() for i, part in enumerate(parts)]
        # result = [english_parts, chinese_parts]
        results_E.append(english_parts)
        results_C.append(chinese_parts)

    # 打印结果
    return results_E, results_C

def Choose_color(text):
    if text == 'black':
        return RGBColor(0, 0, 0)
    elif text == 'gray':
        return RGBColor(128, 128, 128)
    elif text == 'blue':
        return RGBColor(0, 0, 255)
    elif text == 'green':
        return RGBColor(0, 128, 0)
    elif text == 'orange':
        return RGBColor(255, 165, 0)
    elif text == 'red':
        return RGBColor(255, 0, 0)

def fill_in(text_frame, text, font_size = 24, font_color = 'black', 
            font_name = 'Times New Roman', bold = True, line_spacing = 1, 
            State = "E", Num_pra = 0):
    if Num_pra == 0:
        State = "C"

    if State == "E":
        p = text_frame.add_paragraph()
        p.text = text
        p.font.size = Pt(font_size)
        p.font.name = font_name
        p.font.bold = bold
        p.line_spacing = line_spacing
        p.font.color.rgb = Choose_color(font_color)
    elif State == "C":
        p = text_frame.paragraphs[Num_pra]
        q = p.add_run()
        q.text = text
        q.font.size = Pt(font_size)
        q.font.name = font_name
        q.font.bold = bold
        q.line_spacing = line_spacing
        q.font.color.rgb = Choose_color(font_color)

# 指定要读取的天序号
day = input("我宝要听写第几天的内容呢？")
# day = 70

# 指定文件路径
file_path = f'./files/day{day}.xlsx'

# 读取Excel文件
df = pd.read_excel(file_path)
# 查看数据的前5行
print("好的鸭，day " + day + " 的前五行是：\n ")
print(df.head())

# 确定要抽查的单词，以及他们的序号
print("首先是单词背诵部分！")
input2 = input("请输入要抽查英文的单词序号，用空格隔开：")
print("好的鸭")
input5 = input("请输入要抽查英文的短语序号，用空格隔开：")
print("好的鸭")
input1 = input("请输入要抽查中文的单词序号，用空格隔开：")
print("好的鸭")
input4 = input("请输入要抽查中文的短语序号，用空格隔开：")
print("好的鸭")
input3 = input("请输入要抽查单词的变形序号，用空格隔开：")
print("好！我已经知悉，正在做PPT当中！请稍后！")

input1_ = input1.split(' ')
input2_ = input2.split(' ')
input3_ = input3.split(' ')
input4_ = input4.split(' ')
input5_ = input5.split(' ')

# 过滤掉空字符串
input1_ = [char.strip() for char in input1.split(' ') if char.strip()]
input2_ = [char.strip() for char in input2.split(' ') if char.strip()]
input3_ = [char.strip() for char in input3.split(' ') if char.strip()]
input4_ = [char.strip() for char in input4.split(' ') if char.strip()]
input5_ = [char.strip() for char in input5.split(' ') if char.strip()]

recite_E_input = [int(char) for char in input1_]
recite_C_input = [int(char) for char in input2_]
trans_input = [int(char) for char in input3_]
phrase_E_input = [int(char) for char in input4_]
phrase_C_input = [int(char) for char in input5_]

recite_E_indices = [(x-1) % 50 for x in recite_E_input]
recite_C_indices = [(x-1) % 50 for x in recite_C_input]
trans_indices = [x-1 for x in trans_input]
phrase_E_indices = [x-1 for x in phrase_E_input]
phrase_C_indices = [x-1 for x in phrase_C_input]
recite_indices = recite_E_indices + recite_C_indices
phrase_indices = phrase_E_indices + phrase_C_indices

RECITE = df.loc[:,"r"].tolist() 
TRANS = df.loc[:,"t"].tolist()
PHRASE = df.loc[:,"p"].tolist()
RECITE = [item for item in RECITE if pd.notna(item)]
TRANS = [item for item in TRANS if pd.notna(item)]
PHRASE = [item for item in PHRASE if pd.notna(item)]

recite_E_all, recite_C_all = parse_recite(RECITE)
phrase_E_all, phrase_C_all = parse_phrase(PHRASE)
trans_E_all, trans_C_all = sparse_transformation(TRANS)

recite_num = len(RECITE)
phrase_num = len(PHRASE)
trans_num = len(TRANS)

prs = Presentation()
prs.slide_width = Cm(33.85)
prs.slide_height = Cm(19.02)
bullet_slide_layout = prs.slide_layouts[6]    # 空白版式

###第一页PPT###
slide1 = prs.slides.add_slide(bullet_slide_layout)
tf_box_1 = slide1.shapes.add_textbox(left=Cm(0), top=Cm(0), width=Cm(33.85), height=Cm(19))
tf1 = tf_box_1.text_frame
tf1.word_wrap = True
uni0 = 0x2460
phr_index = 0
_attach = '__________'

# 考单词的英文
for i in recite_C_indices:
    s = recite_C_all[i]
    uni_ = uni0 + phr_index
    text = _attach
    fill_in(tf1, text, State='E', Num_pra = phr_index)
    text = "        " + s
    fill_in(tf1, text, State='C', Num_pra = phr_index)
    phr_index += 1
# 考短语的英文
for i in phrase_C_indices:
    s = phrase_C_all[i]
    uni_ = uni0 + phr_index
    text = _attach
    fill_in(tf1,text, State = 'E', Num_pra= phr_index)
    text = "        " + s
    fill_in(tf1,text, State='C', Num_pra = phr_index)
    phr_index += 1
# 考单词的中文
for i in recite_E_indices:
    s = recite_E_all[i]
    uni_ = uni0 + phr_index
    text = s
    fill_in(tf1,text,State='E', Num_pra = phr_index)
    text = "        " + _attach
    fill_in(tf1,text, State='C', Num_pra = phr_index)
    phr_index += 1
# 考短语的中文
for i in phrase_E_indices:
    s = phrase_E_all[i]
    uni_ = uni0 + phr_index
    text = s
    fill_in(tf1, text, State='E', Num_pra = phr_index)
    text = "        " + _attach
    fill_in(tf1,text, State='C', Num_pra = phr_index)
    phr_index += 1
    
for i in trans_indices:
    text = ""
    s_ = trans_C_all[i]
    uni_ = uni0 + phr_index
    # text = chr(uni_)
    fill_in(tf1, text, Num_pra = phr_index)
    for j in s_:
        text = _attach
        fill_in(tf1,text = text,State='C', Num_pra = phr_index)
        text = "     "+ j + "    "
        fill_in(tf1,text = text,State='C', Num_pra = phr_index)
    phr_index += 1

tf_box_3 = slide1.shapes.add_textbox(left=Cm(30.67), top=Cm(17.92), width=Cm(3.18), height=Cm(1.11))
tf3 = tf_box_3.text_frame
tf3.paragraphs[0].text = f"day{day}"
tf3.paragraphs[0].font.name = 'Times New Roman'
tf3.paragraphs[0].font.bold = True
tf3.paragraphs[0].font.size = Pt(24)


###第二页PPT###
slide2 = prs.slides.add_slide(bullet_slide_layout)
tf_box_2 = slide2.shapes.add_textbox(left=Cm(0), top=Cm(0), width=Cm(33.85), height=Cm(19))
tf2 = tf_box_2.text_frame
tf2.word_wrap = True
uni0 = 0x2460
phr_index = 0

# 考单词的英文
for i in recite_C_indices:
    s = recite_C_all[i]
    ss = recite_E_all[i]
    uni_ = uni0 + phr_index
    text = ss
    fill_in(tf2, text, State='E', font_color = 'red', Num_pra = phr_index)
    text = "       " + s
    fill_in(tf2, text, State='C', font_color = 'gray', Num_pra = phr_index)
    phr_index += 1

# 考短语的英文
for i in phrase_C_indices:
    s = phrase_C_all[i]
    ss = phrase_E_all[i]
    uni_ = uni0 + phr_index
    text = ss
    fill_in(tf2, text, State='E', font_color = 'red', Num_pra = phr_index )
    text = "        " + s
    fill_in(tf2, text, State='C', font_color = 'gray', Num_pra = phr_index )
    phr_index += 1

# 考单词的中文
for i in recite_E_indices:
    s = recite_E_all[i]
    ss = recite_C_all[i]
    uni_ = uni0 + phr_index
    text = s
    fill_in(tf2, text, State='E', font_color = 'red', Num_pra = phr_index )
    text = "        " + ss
    fill_in(tf2, text, State= 'C', font_color = 'gray', Num_pra = phr_index )
    phr_index += 1
# 考短语的中文
for i in phrase_E_indices:
    s = phrase_E_all[i]
    ss = phrase_C_all[i]
    uni_ = uni0 + phr_index
    text = s
    fill_in(tf2, text, State='E', font_color = 'red', Num_pra = phr_index )
    text = "        " + ss 
    fill_in(tf2,text, State= 'C', font_color = 'gray', Num_pra = phr_index )
    phr_index += 1
    
for i in trans_indices:
    text = ""
    s_ = trans_C_all[i]
    ss_ = trans_E_all[i]
    uni_ = uni0 + phr_index
    # text = chr(uni_)
    fill_in(tf2, text, State='E', font_color = 'red', Num_pra = phr_index )
    for j,t in zip(s_,ss_):
        text = t
        fill_in(tf2,text, State='C', font_color = 'red', Num_pra = phr_index)
        text = "    " + j + "    "
        fill_in(tf2,text, State='C', font_color = 'gray', Num_pra = phr_index)
    phr_index += 1

tf_box_4 = slide2.shapes.add_textbox(left=Cm(30.67), top=Cm(17.92), width=Cm(3.18), height=Cm(1.11))
tf4 = tf_box_4.text_frame
tf4.paragraphs[0].text = f"day{day}"
tf4.paragraphs[0].font.name = 'Times New Roman'
tf4.paragraphs[0].font.bold = True
tf4.paragraphs[0].font.size = Pt(24)

prs.save(f"./day{day}_dictation.pptx")
print("♥♥ ♥♥ ♥♥ ♥♥ ♥♥ ♥♥ 做好了，爱你一万年！♥♥ ♥♥ ♥♥ ♥♥ ♥♥ ♥♥")
input("按任意键退出程序")